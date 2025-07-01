from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth import authenticate, login
from django.contrib.auth.models import User, Group
from django.contrib import messages
from django.http import JsonResponse, HttpResponse, FileResponse
from django.db.models import Q, Sum, Min, Max, FloatField, Count, Value, Case, When, IntegerField, F
from django.db.models.functions import Coalesce
from django.conf import settings
from django.core.cache import cache
from django.views.decorators.cache import cache_page
import random
import string
import os
import io
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.utils import get_column_letter
from datetime import datetime, timedelta
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4, landscape
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import Paragraph, SimpleDocTemplate, Table, TableStyle, Image, Spacer
from reportlab.lib import colors
from PIL import Image as PILImage, ImageDraw, ImageFont
import tempfile
import subprocess
import zipfile
from django.core.mail import send_mail
from django.contrib.auth.decorators import login_required, user_passes_test
from django.utils import timezone
from dateutil import parser as dateparser
from django.core.paginator import Paginator
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import sys
import platform
import time
import re

from .models import Scanner, Event, EventParticipant, PurgeSettings, NotificationLog

# Проверка доступа (только тимлидеры и админы)
def is_team_leader_or_admin(user):
    return user.is_authenticated and (user.is_team_leader or user.is_staff)

def team_leader_required(view_func):
    return user_passes_test(is_team_leader_or_admin)(view_func)

# Проверка прав администратора
def is_admin(user):
    return user.is_authenticated and user.is_staff

def admin_required(view_func):
    return user_passes_test(is_admin)(view_func)

# Авторизация
def login_request(request):
    if request.method == 'POST':
        name = request.POST.get('name')
        email = request.POST.get('email')
        try:
            user = User.objects.get(email=email)
        except User.DoesNotExist:
            return render(request, 'core/login.html', {'error': 'Только для тимлидеров'})
        code = str(random.randint(100000, 999999))
        request.session['auth_code'] = code
        request.session['auth_email'] = email
        send_mail(
            'Код для входа',
            f'Ваш код: {code}',
            settings.DEFAULT_FROM_EMAIL,
            [email],
            fail_silently=False,
        )
        return redirect('code_verify')
    return render(request, 'core/login.html')

def code_verify(request):
    if request.method == 'POST':
        code = request.POST.get('code')
        if code == request.session.get('auth_code'):
            email = request.session.get('auth_email')
            user = User.objects.get(email=email)
            login(request, user)
            return redirect('events')
        return render(request, 'core/code_verify.html', {'error': 'Неверный код'})
    return render(request, 'core/code_verify.html')

@team_leader_required
def events_list(request):
    events = Event.objects.all().order_by('-date')
    # Поиск по названию
    q = request.GET.get('q', '').strip()
    if q:
        events = events.filter(name__icontains=q)
    # Фильтр по дате
    date_from = request.GET.get('date_from')
    date_to = request.GET.get('date_to')
    if date_from:
        events = events.filter(date__gte=date_from)
    if date_to:
        events = events.filter(date__lte=date_to)
    # Фильтр по тимлидеру
    leader = request.GET.get('leader')
    if leader:
        events = events.filter(created_by__first_name__icontains=leader) | events.filter(created_by__last_name__icontains=leader)
    # Пагинация
    paginator = Paginator(events, 20)
    page_number = request.GET.get('page')
    page_obj = paginator.get_page(page_number)
    return render(request, 'core/events_list.html', {
        'events': page_obj.object_list,
        'page_obj': page_obj,
        'q': q,
        'date_from': date_from,
        'date_to': date_to,
        'leader': leader,
    })

@team_leader_required
def event_create(request):
    scanners = Scanner.objects.all()
    if request.method == 'POST':
        name = request.POST.get('name')
        date = request.POST.get('date')
        start_date = request.POST.get('start_date')
        end_date = request.POST.get('end_date')
        max_scanners = request.POST.get('max_scanners', 10)
        location = request.POST.get('location', '')
        
        # Если указан период, используем его, иначе обычную дату
        if start_date and end_date:
            event = Event.objects.create(
                name=name,
                date=start_date,  # для обратной совместимости
                start_date=start_date,
                end_date=end_date,
                max_scanners=max_scanners,
                location=location,
                description='',
                created_by=request.user
            )
        else:
            event = Event.objects.create(
                name=name,
                date=date,
                start_date=date if date else None,
                end_date=date if date else None,
                max_scanners=max_scanners,
                location=location,
                description='',
                created_by=request.user
            )
        return redirect('event_edit', event_id=event.id)
    return render(request, 'core/event_create.html', {'scanners': scanners})

@team_leader_required
def event_edit(request, event_id):
    event = Event.objects.get(id=event_id)
    all_scanners = Scanner.objects.all()
    participants = EventParticipant.objects.filter(event=event).select_related('volunteer')
    
    if request.method == 'POST':
        if request.user.is_staff or event.created_by == request.user:
            if 'add_volunteer' in request.POST:
                # Проверяем, есть ли множественные ID сканеров (через запятую)
                volunteer_ids = request.POST.getlist('volunteer_ids[]', [])
                # Если нет множественных ID, проверяем одиночный ID
                if not volunteer_ids:
                    volunteer_id = request.POST.get('volunteer_id')
                    if volunteer_id:
                        volunteer_ids = [volunteer_id]
                
                # Проверяем, есть ли строка с именами для поиска
                volunteer_names = request.POST.get('volunteer_names', '').strip()
                if volunteer_names:
                    # Разбиваем строку на имена и ищем сканеров
                    names = volunteer_names.split()
                    if names:
                        # Создаем запрос для поиска сканеров по именам
                        query = Q()
                        for name in names:
                            if name:
                                query |= Q(first_name__icontains=name) | Q(last_name__icontains=name)
                        
                        # Получаем ID найденных сканеров
                        found_scanners = Scanner.objects.filter(query).values_list('id', flat=True)
                        volunteer_ids.extend([str(id) for id in found_scanners])
                
                # Убираем дубликаты
                volunteer_ids = list(set(volunteer_ids))
                
                # Добавляем сканеров
                current_participants_count = participants.count()
                added_count = 0
                not_added_count = 0
                
                for volunteer_id in volunteer_ids:
                    # Проверяем, не превышен ли лимит
                    if current_participants_count + added_count >= event.max_scanners:
                        messages.error(request, f'Превышено максимальное количество волонтеров ({event.max_scanners})')
                        break
                    
                    # Проверяем, не добавлен ли уже этот сканер
                    if not participants.filter(volunteer_id=volunteer_id).exists():
                        try:
                            volunteer = Scanner.objects.get(id=volunteer_id)
                            EventParticipant.objects.create(event=event, volunteer=volunteer)
                            added_count += 1
                        except Scanner.DoesNotExist:
                            not_added_count += 1
                    else:
                        not_added_count += 1
                
                if added_count > 0:
                    messages.success(request, f'Добавлено {added_count} сканеров')
                if not_added_count > 0:
                    messages.warning(request, f'{not_added_count} сканеров не были добавлены (уже участвуют или не найдены)')
                
                return redirect('event_edit', event_id=event.id)
            if 'remove_participant' in request.POST and request.user.is_staff:
                participant_id = request.POST.get('participant_id')
                try:
                    participant = EventParticipant.objects.get(id=participant_id, event=event)
                    participant.delete()
                    messages.success(request, f'Участник удален из мероприятия')
                except EventParticipant.DoesNotExist:
                    messages.error(request, 'Участник не найден')
                return redirect('event_edit', event_id=event.id)
            if 'save_participants' in request.POST and request.user.is_staff:
                messages.success(request, 'Изменения сохранены')
                pass
            if 'set_duration' in request.POST and request.user.is_staff:
                try:
                    duration_hours = float(request.POST.get('duration_hours', 0))
                    event.duration_hours = duration_hours
                    event.save()
                    for p in participants:
                        awarded_hours = duration_hours
                        p.hours_awarded = awarded_hours
                        p.save()
                    messages.success(request, f'Продолжительность мероприятия установлена: {duration_hours} часов')
                except ValueError:
                    messages.error(request, 'Неверное значение продолжительности')
        else:
            messages.error(request, 'У вас нет прав для редактирования этого мероприятия')
        clear_scanners_cache()
        return redirect('event_edit', event_id=event.id)
    return render(request, 'core/event_edit.html', {
        'event': event,
        'all_scanners': all_scanners,
        'participants': participants,
        'is_creator': event.created_by == request.user,
        'is_admin': request.user.is_staff
    })

@team_leader_required
def volunteer_search(request):
    q = request.GET.get('q', '').strip()
    event_id = request.GET.get('event_id')
    exclude_ids = list(EventParticipant.objects.filter(event_id=event_id).values_list('volunteer_id', flat=True)) if event_id else []
    
    # Если в запросе несколько имен через пробел, разбиваем их
    search_terms = q.split()
    if search_terms:
        # Создаем пустой Q-объект для OR-условий
        query = Q()
        # Для каждого термина поиска добавляем условия
        for term in search_terms:
            term = term.strip()
            if term:
                # Добавляем условия через OR для каждого термина
                query |= Q(first_name__icontains=term) | Q(last_name__icontains=term) | Q(email__icontains=term)
    
        # Применяем собранный запрос
        scanners = Scanner.objects.filter(query).exclude(id__in=exclude_ids).order_by('last_name', 'first_name')
    else:
        scanners = Scanner.objects.none()
    
    data = [
        {'id': v.id, 'name': f'{v.first_name} {v.last_name}', 'email': v.email}
        for v in scanners
    ]
    
    # Очищаем кеш после добавления нового сканера
    clear_scanners_cache()
    
    return JsonResponse({'results': data})

def home(request):
    if request.user.is_authenticated:
        return redirect('events')
    else:
        return redirect('login')

def export_event_participants(request, event_id):
    event = Event.objects.get(id=event_id)
    participants = EventParticipant.objects.filter(event=event).select_related('volunteer')
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = 'Участники'
    ws.append(['ФИО', 'Email', 'Опоздал (мин)', 'Зачтено (часы:минуты)'])
    for p in participants:
        fio = f'{p.volunteer.first_name} {p.volunteer.last_name}'
        email = p.volunteer.email
        late = p.late_minutes or 0
        hours = int(p.hours_awarded)
        minutes = int(round((p.hours_awarded - hours) * 60))
        time_str = f'{hours}:{minutes:02d}'
        ws.append([fio, email, late, time_str])
    for col in range(1, 5):
        ws.column_dimensions[get_column_letter(col)].width = 22
    response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
    response['Content-Disposition'] = f'attachment; filename=event_{event.id}_participants.xlsx'
    wb.save(response)
    return response



def export_all_events(request):
    events = Event.objects.all().order_by('-date')
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = 'События'
    
    # Стили для заголовков
    header_font = Font(bold=True)
    header_fill = PatternFill(start_color="DDDDDD", end_color="DDDDDD", fill_type="solid")
    
    # Заголовки
    headers = ['Название', 'Дата', 'Тимлидер', 'Треб. волонтёров', 'Длительность (ч)', 'Факт. волонтёров']
    for col_num, header in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col_num, value=header)
        cell.font = header_font
        cell.fill = header_fill
    
    # Данные
    row_num = 2
    for event in events:
        participants_count = EventParticipant.objects.filter(event=event).count()
        ws.cell(row=row_num, column=1, value=event.name)
        ws.cell(row=row_num, column=2, value=event.date)
        ws.cell(row=row_num, column=3, value=f"{event.created_by.first_name} {event.created_by.last_name}")
        ws.cell(row=row_num, column=4, value=event.max_scanners)
        ws.cell(row=row_num, column=5, value=event.duration_hours)
        ws.cell(row=row_num, column=6, value=participants_count)
        row_num += 1
    
    # Автоподбор ширины столбцов
    for col in range(1, len(headers) + 1):
        ws.column_dimensions[get_column_letter(col)].width = 18
    
    response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
    response['Content-Disposition'] = 'attachment; filename=all_events.xlsx'
    wb.save(response)
    return response

def save_pptx_as_png(pptx_data):
    """
    Функция-заглушка для работы с PPTX. На данный момент просто возвращает PPTX как есть,
    так как конвертация в PNG требует дополнительных инструментов.
    """
    # Возвращаем оригинальный PPTX
    return pptx_data.getvalue() if isinstance(pptx_data, io.BytesIO) else pptx_data, False

def convert_pptx_to_png(pptx_path):
    """
    Конвертирует PPTX в PNG используя автоматизацию PowerPoint (для Windows)
    или LibreOffice (для других ОС).
    """
    temp_png_path = pptx_path.replace('.pptx', '.png')
    
    if platform.system() == 'Windows':
        try:
            # Пробуем использовать COM-объект PowerPoint с правильной инициализацией
            import comtypes.client
            import pythoncom
            
            # Инициализация COM
            pythoncom.CoInitialize()
            
            # Создаем экземпляр PowerPoint
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = True
            
            # Полный путь к файлу (абсолютный)
            abs_pptx_path = os.path.abspath(pptx_path)
            presentation = powerpoint.Presentations.Open(abs_pptx_path)
            
            # Сохраняем первый слайд как PNG
            slide = presentation.Slides[0]
            abs_png_path = os.path.abspath(temp_png_path)
            slide.Export(abs_png_path, "PNG", 1024, 768)  # Указываем размер
            
            # Закрываем все
            presentation.Close()
            powerpoint.Quit()
            pythoncom.CoUninitialize()  # Освобождаем COM
            
            if os.path.exists(temp_png_path):
                with open(temp_png_path, 'rb') as f:
                    png_data = f.read()
                os.remove(temp_png_path)  # Удаляем временный файл
                return png_data, True
            
        except Exception as e:
            print(f"Ошибка при конвертации через PowerPoint: {e}")
            try:
                # В случае ошибки, попробуем освободить COM
                pythoncom.CoUninitialize()
            except:
                pass
    
    # Если PowerPoint не сработал, попробуем создать скриншот PPTX напрямую
    try:
        from PIL import Image, ImageDraw, ImageFont
        from pptx import Presentation
        
        # Открываем презентацию для анализа
        prs = Presentation(pptx_path)
        
        # Создаем изображение размером со слайд
        width, height = 1920, 1080  # Стандартный размер слайда 16:9
        img = Image.new("RGB", (width, height), (235, 199, 0))  # Желтый фон как на фото
        draw = ImageDraw.Draw(img)
        
        # Добавляем черные полосы сверху и снизу как на фото (кинолента)
        stripe_height = 50
        for i in range(0, width, 150):
            # Верхняя полоса
            draw.rectangle([(i, 0), (i+100, stripe_height)], fill=(30, 30, 30))
            # Нижняя полоса
            draw.rectangle([(i, height-stripe_height), (i+100, height)], fill=(30, 30, 30))
        
        # Извлекаем данные из PPTX
        # Так как у нас нет прямого доступа к информации, которая использовалась для заполнения,
        # мы будем искать текст в слайдах
        
        name = ""
        hours = ""
        event_name = ""
        event_date = ""
        
        # Проходим по всем слайдам и текстовым блокам
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text = run.text
                            
                            # Здесь очень примерно определяем, что этот текст может быть
                            if len(text) > 0:
                                text = text.strip()
                                if text.isupper() and len(name) == 0 and len(text) > 5:
                                    name = text  # Вероятно это имя (в верхнем регистре)
                                elif text.isdigit() or (len(text) < 3 and text.replace('.', '').isdigit()):
                                    hours = text  # Вероятно это часы
                                elif "20" in text and ("." in text or "/" in text) and len(event_date) == 0:
                                    event_date = text  # Вероятно это дата
                                elif len(text) > 10 and len(event_name) == 0:
                                    event_name = text  # Вероятно это название мероприятия
        
        # Попробуем найти шрифты
        try:
            # В Windows пути к шрифтам
            fonts_dir = os.path.join(os.environ['SystemRoot'], 'Fonts')
            
            # Стандартные шрифты, которые должны быть в большинстве систем
            arial_path = os.path.join(fonts_dir, 'Arial.ttf')
            arialbd_path = os.path.join(fonts_dir, 'Arialbd.ttf')
            impact_path = os.path.join(fonts_dir, 'Impact.ttf')
            
            # Проверяем наличие шрифтов
            if os.path.exists(arialbd_path):
                name_font = ImageFont.truetype(arialbd_path, 60)
            else:
                name_font = ImageFont.load_default()
                
            if os.path.exists(impact_path):
                hours_font = ImageFont.truetype(impact_path, 80)
            else:
                hours_font = ImageFont.load_default()
                
            if os.path.exists(arial_path):
                normal_font = ImageFont.truetype(arial_path, 36)
            else:
                normal_font = ImageFont.load_default()
        except:
            # Если не нашли, используем стандартный
            name_font = ImageFont.load_default()
            hours_font = ImageFont.load_default()
            normal_font = ImageFont.load_default()
        
        # Фон для лучшей видимости
        draw.rectangle([(0, 0), (width, height)], fill="white")
        
        # Рисуем данные на изображении
        try:
            # Попробуем нарисовать билет в левом верхнем углу
            ticket_width, ticket_height = 200, 200
            ticket_x, ticket_y = 80, 150
            
            # Рамка билета
            draw.rectangle(
                [(ticket_x, ticket_y), (ticket_x + ticket_width, ticket_y + ticket_height)],
                outline=(0, 0, 0),
                width=3
            )
            
            # Текст на билете
            small_font = ImageFont.load_default()
            if os.path.exists(arial_path):
                small_font = ImageFont.truetype(arial_path, 20)
                
            draw.text((ticket_x + 10, ticket_y + 10), "TICKET", font=small_font, fill=(0, 0, 0))
            draw.text((ticket_x + 10, ticket_y + 40), "VIP", font=small_font, fill=(0, 0, 0))
            
            # Звезды вокруг билета
            for i in range(8):
                star_x = ticket_x + random.randint(0, ticket_width)
                star_y = ticket_y + random.randint(0, ticket_height)
                star_size = random.randint(10, 20)
                draw.text((star_x, star_y), "★", font=normal_font, fill=(0, 0, 0))
                
        except Exception as e:
            print(f"Ошибка при создании билета: {e}")
        
        # Основной заголовок "CERTIFICAT" большими буквами как на фото
        title = "CERTIFICAT"
        title_font_size = 150
        title_font = hours_font
        if os.path.exists(impact_path):
            try:
                title_font = ImageFont.truetype(impact_path, title_font_size)
            except:
                pass
            
        title_width = title_font.getlength(title) if hasattr(title_font, 'getlength') else len(title) * title_font_size // 2
        title_x = (width - title_width) // 2
        title_y = 100
        
        # Белый текст как на фото
        draw.text((title_x, title_y), title, font=title_font, fill=(255, 255, 255))
        
        # Подзаголовок "OF APPRECIATION"
        subtitle = "OF APPRECIATION"
        subtitle_font_size = 60
        subtitle_font = normal_font
        if os.path.exists(arial_path):
            try:
                subtitle_font = ImageFont.truetype(arial_path, subtitle_font_size)
            except:
                pass
        
        subtitle_width = subtitle_font.getlength(subtitle) if hasattr(subtitle_font, 'getlength') else len(subtitle) * subtitle_font_size // 2
        subtitle_x = width - subtitle_width - 100
        subtitle_y = title_y + title_font_size - 30
        
        draw.text((subtitle_x, subtitle_y), subtitle, font=subtitle_font, fill=(255, 255, 255))
        
        # Текст "This certificate is presented to:"
        pres_text = "This certificate is presented to:"
        pres_font = ImageFont.truetype(sans_italic, size=18 * scale)
        pres_w, pres_h = get_text_size(pres_text, pres_font)
        pres_x = (width - pres_w - 120 * scale)
        pres_y = subtitle_y + subtitle_font_size - 30
        draw.text((pres_x, pres_y), pres_text, font=pres_font, fill=(255,255,255,255))
        
        # Имя участника (большим зеленым шрифтом)
        if name:
            name_font_size = 100
            name_font = hours_font
            if os.path.exists(impact_path):
                try:
                    name_font = ImageFont.truetype(impact_path, name_font_size)
                except:
                    pass
            
            name_width = name_font.getlength(name) if hasattr(name_font, 'getlength') else len(name) * name_font_size // 2
            name_x = (width - name_width) // 2
            name_y = subtitle_y + 120
            
            # Зеленый текст как на фото
            draw.text((name_x, name_y), name, font=name_font, fill=(120, 220, 80))
        
        # Определяем годы для текста благодарности
        years_text = "2024 and 2025"  # По умолчанию
        
        # Если есть список мероприятий, извлекаем годы из них
        if events_list:
            years = set()
            for event in events_list:
                if 'date' in event:
                    try:
                        # Извлекаем год из даты в формате DD.MM.YYYY
                        date_parts = event['date'].split('.')
                        if len(date_parts) == 3:
                            years.add(date_parts[2])  # Год - последняя часть
                    except:
                        pass
            
            # Если нашли годы, формируем строку
            if years:
                sorted_years = sorted(years)
                if len(sorted_years) == 1:
                    years_text = sorted_years[0]
                elif len(sorted_years) == 2:
                    years_text = f"{sorted_years[0]} and {sorted_years[1]}"
                else:
                    # Для трех и более лет используем перечисление с запятыми и "and" перед последним
                    years_text = ", ".join(sorted_years[:-1]) + " and " + sorted_years[-1]
        
        # Если передан период, извлекаем годы из него
        elif period:
            try:
                # Предполагаем, что период имеет формат "DD.MM.YYYY - DD.MM.YYYY"
                period_parts = period.split(' - ')
                if len(period_parts) == 2:
                    start_year = period_parts[0].split('.')[-1]
                    end_year = period_parts[1].split('.')[-1]
                    
                    if start_year == end_year:
                        years_text = start_year
                    else:
                        years_text = f"{start_year} and {end_year}"
            except:
                pass
        
        # Если передана одиночная дата мероприятия
        elif event_date:
            try:
                # Предполагаем формат даты "DD.MM.YYYY" или объект datetime
                if isinstance(event_date, str):
                    year = event_date.split('.')[-1]
                    years_text = year
                else:
                    years_text = str(event_date.year)
            except:
                pass

        # Блок благодарности с динамическими годами
        lines = [
            "We, the Ticketon company, would like to sincerely express our gratitude and",
            "appreciation towards your incredible work and support in organizing",
            f"our events in {years_text}. You played important role in organization of",
            "each event. We hope to see you again in upcoming events!"
        ]
        main_font = ImageFont.truetype(sans_italic, size=16 * scale)
        main_width = 600 * scale
        main_x = width - main_width - 40 * scale
        main_y = name_y + name_h + 40 * scale
        line_h = main_font.getbbox('Ag')[3] - main_font.getbbox('Ag')[1] + 8 * scale
        for i, line in enumerate(lines):
            words = line.split()
            if len(words) == 1:
                draw.text((main_x, main_y + i * line_h), line, font=main_font, fill=(255,255,255,255))
                continue
            total_w = sum(get_text_size(w, main_font)[0] for w in words)
            space_w = int((main_width - total_w) / max(1, (len(words) - 1) * 2))
            x = main_x
            for j, word in enumerate(words):
                draw.text((x, main_y + i * line_h), word, font=main_font, fill=(255,255,255,255))
                w, _ = get_text_size(word, main_font)
                x += w + space_w

        # Рисуем стрелку с остриём только слева, справа — ровно
        arrow_w, arrow_height = 540 * scale, 110 * scale
        arrow = PILImage.new('RGBA', (arrow_w, arrow_height), (0,0,0,0))
        adraw = ImageDraw.Draw(arrow)
        points = [
            (0, arrow_height//2), (40 * scale, 0), (arrow_w, 0), (arrow_w, arrow_height), (40 * scale, arrow_height)
        ]
        adraw.polygon(points, fill=(76,175,80,255))

        margin = 18 * scale
        gap = 14 * scale
        gap_hours = 6 * scale
        section_w = (arrow_w - 2 * margin - 2 * gap) // 3
        center_y = arrow_height // 2

        # Шрифты для стрелки
        impact_25 = ImageFont.truetype(impact_path, size=25 * scale)
        impact_19_9 = ImageFont.truetype(impact_path, size=int(19.9 * scale))
        impact_18 = ImageFont.truetype(impact_path, size=18 * scale)

        # Левая секция: часы
        hours_text = f"{int(round(hours)):02d}"
        hw, hh = get_text_size(hours_text, impact_25)
        hlabel = "hours"
        hlw, hlh = get_text_size(hlabel, impact_18)
        line_w = section_w * 0.8
        total_h = hh + gap + hlh
        base_y = center_y - total_h // 2 - 10 * scale
        hours_x = margin + (section_w - hw) // 2
        hours_y = base_y
        adraw.text((hours_x, hours_y), hours_text, font=impact_25, fill=(255,255,255,255))
        line_y = hours_y + hh + gap
        line_x1 = margin + (section_w - line_w) // 2
        line_x2 = line_x1 + line_w
        adraw.line([(line_x1, line_y), (line_x2, line_y)], fill=(255,255,255,255), width=4 * scale)
        hlabel_x = margin + (section_w - hlw) // 2
        hlabel_y = line_y + gap_hours
        adraw.text((hlabel_x, hlabel_y), hlabel, font=impact_18, fill=(255,255,255,255))

        # Центральная секция: печать (штамп чуть выше, овальный)
        if os.path.exists(stamp_path):
            stamp = PILImage.open(stamp_path).convert('RGBA')
            stamp_w = int(section_w)
            stamp_h = int(arrow_height * 0.95)
            stamp = stamp.resize((stamp_w, stamp_h), PILImage.LANCZOS)
            stamp_x = margin + section_w + gap
            stamp_y = center_y - stamp_h // 2
            arrow.paste(stamp, (stamp_x, stamp_y), stamp)

        # Правая секция: крупный текст, увеличиваем размер имени директора и слова "director"
        sign_text = "Torgunakova V. K."
        dir_text = "director"
        max_width = section_w - 2 * int(5 * scale)  # Уменьшаем внутренний отступ для большего текста
        
        # Увеличиваем базовый размер шрифта для имени директора и слова "director"
        base_sign_font_size = int(25 * 1.5 * scale)  # Было 20.9 * 1.5
        base_dir_font_size = int(22 * 1.5 * scale)   # Было 19 * 1.5
        
        min_font_size = int(15 * scale)  # Увеличиваем минимальный размер шрифта
        sign_font_size = base_sign_font_size
        dir_font_size = base_dir_font_size
        
        while True:
            sign_font = ImageFont.truetype(impact_path, size=sign_font_size)
            dir_font = ImageFont.truetype(impact_path, size=dir_font_size)
            sign_w, sign_h = get_text_size(sign_text, sign_font)
            dir_w, dir_h = get_text_size(dir_text, dir_font)
            if sign_w <= max_width and dir_w <= max_width:
                break
            sign_font_size -= 2
            dir_font_size -= 2
            if sign_font_size < min_font_size or dir_font_size < min_font_size:
                sign_font_size = dir_font_size = min_font_size
                sign_font = ImageFont.truetype(impact_path, size=sign_font_size)
                dir_font = ImageFont.truetype(impact_path, size=dir_font_size)
                sign_w, sign_h = get_text_size(sign_text, sign_font)
                dir_w, dir_h = get_text_size(dir_text, dir_font)
                break
        
        right_section_x = margin + 2 * section_w + 2 * gap
        right_section_y = margin
        sign_x = right_section_x + (section_w - sign_w) // 2
        sign_y = right_section_y + 10 * scale
        dir_x = right_section_x + (section_w - dir_w) // 2
        dir_y = sign_y + sign_h + int(0.2 * sign_h)  # Уменьшаем расстояние между именем и должностью
        
        # Рисуем имя и должность директора
        adraw.text((sign_x, sign_y), sign_text, font=sign_font, fill=(0,0,0,255))
        adraw.text((dir_x, dir_y), dir_text, font=dir_font, fill=(255,255,255,255))

        arrow_x = width - arrow_w
        arrow_y = height - arrow_height - 60 * scale
        img.paste(arrow, (arrow_x, arrow_y), arrow)

        temp_dir = tempfile.mkdtemp()
        temp_img_path = os.path.join(temp_dir, 'cert.png')
        img.save(temp_img_path, 'PNG')

        # Использование ReportLab для создания PDF без артефактов
        pdf_width, pdf_height = 1123, 794  # Размеры PDF в точках
        temp_pdf_path = os.path.join(temp_dir, 'certificate.pdf')
        
        # Создаем PDF с чистым черным фоном, без артефактов
        c = canvas.Canvas(temp_pdf_path, pagesize=(pdf_width, pdf_height))
        
        # Заливаем весь PDF черным цветом (без границ)
        c.setFillColor((0, 0, 0))
        c.rect(0, 0, pdf_width, pdf_height, fill=1, stroke=0)
        
        # Добавляем изображение сертификата (без верхней черной линии)
        c.drawImage(temp_img_path, 0, 0, width=pdf_width, height=pdf_height)
        
        c.save()

        with open(temp_pdf_path, 'rb') as f:
            pdf_data = f.read()
        
        # Очистка временных файлов
        os.remove(temp_img_path)
        os.remove(temp_pdf_path)
        os.rmdir(temp_dir)
        
        return pdf_data
    except Exception as e:
        print(f"Ошибка при создании PNG: {e}")
    
    # Пробуем через LibreOffice в последнюю очередь
    try:
        if platform.system() == 'Windows':
            # Если у нас есть LibreOffice, попробуем его использовать
            libreoffice_paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
            ]
            
            soffice_path = None
            for path in libreoffice_paths:
                if os.path.exists(path):
                    soffice_path = path
                    break
                
            if soffice_path:
                # Полный путь к файлу и директории
                abs_pptx_path = os.path.abspath(pptx_path)
                output_dir = os.path.dirname(abs_pptx_path)
                
                cmd = [
                    soffice_path, '--headless', 
                    '--convert-to', 'png', '--outdir', 
                    output_dir, abs_pptx_path
                ]
                subprocess.run(cmd, timeout=30, check=True)
            else:
                return None, False
        else:
            cmd = ['soffice', '--headless', '--convert-to', 'png', 
                  '--outdir', os.path.dirname(pptx_path), pptx_path]
            subprocess.run(cmd, timeout=30, check=True)
        
        # LibreOffice создаст файл с тем же именем, но с расширением .png
        if os.path.exists(temp_png_path):
            with open(temp_png_path, 'rb') as f:
                png_data = f.read()
            os.remove(temp_png_path)  # Удаляем временный файл
            return png_data, True
    except Exception as e:
        print(f"Ошибка при конвертации через LibreOffice: {e}")
    
    # Если все предыдущие попытки не удались, просто используем PPTX
    with open(pptx_path, 'rb') as f:
        pptx_data = f.read()
    return pptx_data, False

def get_certificate_from_template(name, hours, event_name=None, event_date=None, leader_name=None, period=None, events_list=None):
    """
    Создает PPTX на основе шаблона и заполняет данными
    """
    # Путь к шаблону
    template_path = os.path.join(settings.BASE_DIR, 'static', 'templates', 'certificate_template.pptx')
    
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Файл шаблона не найден по пути {template_path}")
    
    # Открываем шаблон
    try:
        prs = Presentation(template_path)
    except Exception as e:
        print(f"Ошибка при открытии шаблона: {e}")
        raise
    
    # Формируем строку списка мероприятий
    events_text = ""
    if events_list:
        for event in events_list:
            events_text += f"• {event['name']} ({event['date']}): {round(event['hours'])} ч.\n"
    
    # Заполняем данные
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Заменяем плейсхолдеры на реальные данные
                        text = run.text
                        
                        # Заменяем NAME на имя сканера
                        if "NAME" in text:
                            run.text = text.replace("NAME", name)
                        
                        # Заменяем 00 на количество часов с форматированием
                        if "00" in text:
                            run.text = text.replace("00", str(round(hours)))
                            run.font.name = "Impact"
                            run.font.size = Pt(25)
                        
                        # Заменяем EVENT_NAME на название мероприятия
                        if "EVENT_NAME" in text or "{{EVENT}}" in text:
                            if event_name:
                                run.text = text.replace("EVENT_NAME", event_name).replace("{{EVENT}}", event_name)
                        
                        # Заменяем EVENT_DATE на дату мероприятия
                        if "EVENT_DATE" in text or "{{DATE}}" in text:
                            if event_date:
                                run.text = text.replace("EVENT_DATE", event_date).replace("{{DATE}}", event_date)
                        
                        # Заменяем LEADER_NAME на имя лидера
                        if "LEADER_NAME" in text or "{{LEADER}}" in text:
                            if leader_name:
                                run.text = text.replace("LEADER_NAME", leader_name).replace("{{LEADER}}", leader_name)
                        
                        # Заменяем PERIOD на период
                        if "PERIOD" in text:
                            if period:
                                run.text = text.replace("PERIOD", period)
                        
                        # Заменяем EVENT_LIST на список мероприятий
                        if "EVENT_LIST" in text:
                            run.text = text.replace("EVENT_LIST", events_text)
    
    # Сохраняем во временный файл
    temp_dir = tempfile.mkdtemp()
    temp_pptx = os.path.join(temp_dir, "certificate.pptx")
    
    try:
        prs.save(temp_pptx)
        
        # Возвращаем путь к временному файлу
        return temp_pptx, temp_dir
    except Exception as e:
        print(f"Ошибка при сохранении PPTX: {e}")
        # Очищаем временную директорию при ошибке
        try:
            os.remove(temp_pptx)
            os.rmdir(temp_dir)
        except:
            pass
        raise

@team_leader_required
def generate_certificate(request, participant_id):
    try:
        participant = EventParticipant.objects.get(id=participant_id)
        event = participant.event
        scanner = participant.volunteer
        full_name = f"{scanner.first_name} {scanner.last_name}".upper()
        hours = round(participant.hours_awarded)
        event_name = event.name
        event_date = event.date.strftime("%d.%m.%Y")
        leader_name = f"{event.created_by.first_name} {event.created_by.last_name}" if event.created_by else None
        file_data = create_certificate_pdf(full_name, hours, event_name, event_date, leader_name)
        participant.hours_awarded = 0
        participant.save()
        filename = f"certificate_{scanner.last_name}_{event.name}.pdf"
        
        # Всегда возвращаем PDF напрямую
        response = HttpResponse(file_data, content_type='application/pdf')
        response['Content-Disposition'] = f'attachment; filename="{filename}"'
        return response
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=400)

@team_leader_required
def generate_all_certificates(request, event_id):
    """
    Генерирует сертификаты для всех участников мероприятия
    """
    try:
        event = Event.objects.get(id=event_id)
        participants = EventParticipant.objects.filter(event=event)
        
        # Создаем zip-архив с сертификатами
        temp_zip = io.BytesIO()
        
        with zipfile.ZipFile(temp_zip, 'w') as zipf:
            for participant in participants:
                scanner = participant.volunteer
                
                # Полное имя сканера для замены (в верхнем регистре)
                full_name = f"{scanner.first_name} {scanner.last_name}".upper()
                
                # Подготавливаем данные для сертификата
                event_name = event.name
                event_date = event.date.strftime("%d.%m.%Y")
                leader_name = f"{event.created_by.first_name} {event.created_by.last_name}" if event.created_by else None
                hours = participant.hours_awarded
                
                # Создаем PDF сертификат напрямую
                pdf_data = create_certificate_pdf(
                    name=full_name, 
                    hours=hours,
                    event_name=event_name,
                    event_date=event_date,
                    leader_name=leader_name
                )
                
                # Обнуляем часы сканера при получении сертификата
                participant.hours_awarded = 0
                participant.save()
                
                # Добавляем в архив
                filename = f"certificate_{scanner.last_name}_{scanner.first_name}.pdf"
                zipf.writestr(filename, pdf_data)
        
        temp_zip.seek(0)
        
        # Отдаем архив для скачивания
        response = FileResponse(
            temp_zip,
            as_attachment=True,
            filename=f"certificates_{event.name}.zip"
        )
        return response
        
    except Exception as e:
        import traceback
        print(f"Ошибка при генерации сертификатов: {e}")
        print(traceback.format_exc())
        return JsonResponse({"error": str(e)}, status=400)

@team_leader_required
def generate_scanner_certificate(request, scanner_id):
    try:
        scanner = get_object_or_404(Scanner, id=scanner_id)
        participations = EventParticipant.objects.filter(volunteer=scanner).select_related('event')
        if not participations.exists():
            messages.error(request, "Сканер еще не участвовал на мероприятиях")
            return redirect('scanner_certificates')
        
        # Получаем текущие часы, доступные для сертификата
        current_hours = participations.aggregate(total_hours=Sum('hours_awarded', output_field=FloatField()))['total_hours'] or 0.0
        
        if current_hours == 0:
            messages.error(request, "У сканера нет доступных часов для получения сертификата")
            return redirect('scanner_certificates')
        
        # Получаем диапазон дат мероприятий
        date_range = participations.aggregate(first_event=Min('event__date'), last_event=Max('event__date'))
        first_date = date_range['first_event']
        last_date = date_range['last_event']
        period_text = f"{first_date.strftime('%d.%m.%Y')} - {last_date.strftime('%d.%m.%Y')}"
        
        # Формируем список мероприятий для сертификата
        events_list = [{
            'name': p.event.name,
            'date': p.event.date.strftime("%d.%m.%Y"),
            'hours': p.hours_awarded
        } for p in participations]
        
        full_name = f"{scanner.first_name} {scanner.last_name}".upper()
        hours = round(current_hours)
        
        # Обновляем общее количество часов, полученных в сертификатах
        scanner.total_certificate_hours += current_hours
        scanner.save()
        
        # Создаем сертификат
        file_data = create_certificate_pdf(full_name, hours, period=period_text, events_list=events_list)
        
        # Обнуляем часы сканера при получении сертификата
        for participant in participations:
            participant.hours_awarded = 0
            participant.save()
        
        filename = f"certificate_{scanner.last_name}_{scanner.first_name}.pdf"
        
        # Всегда возвращаем PDF напрямую
        response = HttpResponse(file_data, content_type='application/pdf')
        response['Content-Disposition'] = f'attachment; filename="{filename}"'
        return response
    except Exception as e:
        import traceback
        trace = traceback.format_exc()
        messages.error(request, f"Ошибка при создании сертификата: {str(e)}")
        return redirect('scanner_certificates')

@team_leader_required
def generate_all_scanner_certificates(request):
    """
    Генерирует сертификаты для всех сканеров
    """
    try:
        # Получаем всех сканеров, которые участвовали хотя бы в одном мероприятии
        scanners = Scanner.objects.filter(id__in=EventParticipant.objects.values_list('volunteer', flat=True).distinct())
        
        # Создаем zip-архив с сертификатами
        temp_zip = io.BytesIO()
        
        with zipfile.ZipFile(temp_zip, 'w') as zipf:
            for scanner in scanners:
                # Получаем все мероприятия сканера
                participations = EventParticipant.objects.filter(volunteer=scanner).select_related('event')
                
                if not participations.exists():
                    continue
                
                # Получаем суммарные часы
                total_hours = participations.aggregate(total_hours=Coalesce(Sum('hours_awarded'), 0))['total_hours']
                
                if total_hours == 0:
                    continue
                
                # Получаем даты первого и последнего мероприятия
                date_range = participations.aggregate(
                    first_event=Min('event__date'),
                    last_event=Max('event__date')
                )
                first_date = date_range['first_event']
                last_date = date_range['last_event']
                
                # Формируем список мероприятий
                events_list = []
                for p in participations:
                    events_list.append({
                        'name': p.event.name,
                        'date': p.event.date.strftime("%d.%m.%Y"),
                        'hours': p.hours_awarded
                    })
                
                # Полное имя сканера (в верхнем регистре)
                full_name = f"{scanner.first_name} {scanner.last_name}".upper()
                
                # Период участия
                period_text = f"{first_date.strftime('%d.%m.%Y')} - {last_date.strftime('%d.%m.%Y')}"
                
                # Создаем PDF сертификат напрямую
                pdf_data = create_certificate_pdf(
                    name=full_name,
                    hours=total_hours,
                    period=period_text,
                    events_list=events_list
                )
                
                # Обнуляем часы сканера при получении сертификата
                for p in participations:
                    p.hours_awarded = 0
                    p.save()
                
                # Добавляем в архив
                filename = f"certificate_{scanner.last_name}_{scanner.first_name}.pdf"
                zipf.writestr(filename, pdf_data)
        
        temp_zip.seek(0)
        
        # Отдаем архив для скачивания
        response = FileResponse(
            temp_zip,
            as_attachment=True,
            filename=f"all_scanner_certificates.zip"
        )
        return response
        
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=400)

@team_leader_required
def scanner_certificates(request):
    """
    Страница поиска сканеров для генерации сертификатов
    """
    return render(request, 'core/scanner_certificate.html')

@team_leader_required
def scanner_events_api(request, scanner_id):
    try:
        scanner = get_object_or_404(Scanner, id=scanner_id)
        participations = EventParticipant.objects.filter(volunteer=scanner).select_related('event')
        
        if not participations.exists():
            messages.info(request, "У сканера нет мероприятий")
            return JsonResponse({"events": [], "total_hours": 0})
        
        # Текущие доступные часы (которые можно использовать для сертификата)
        current_hours = participations.aggregate(total_hours=Sum('hours_awarded', output_field=FloatField()))['total_hours'] or 0.0
        
        # Общее количество часов (включая уже полученные в сертификатах)
        total_all_hours = current_hours + scanner.total_certificate_hours
        
        date_range = participations.aggregate(first_event=Min('event__date'), last_event=Max('event__date'))
        first_date = date_range['first_event']
        last_date = date_range['last_event']
        
        events = [{
            'name': p.event.name,
            'date': p.event.date.strftime("%d.%m.%Y"),
            'hours': p.hours_awarded
        } for p in participations]
        
        return JsonResponse({
            "events": events,
            "current_hours": current_hours,  # Текущие доступные часы
            "total_certificate_hours": scanner.total_certificate_hours,  # Часы, полученные в сертификатах
            "total_all_hours": total_all_hours,  # Общее количество часов
            "first_date": first_date.strftime("%d.%m.%Y") if first_date else None,
            "last_date": last_date.strftime("%d.%m.%Y") if last_date else None
        })
    except Exception as e:
        messages.error(request, f"Ошибка при получении данных: {str(e)}")
        return JsonResponse({"error": str(e)}, status=400)

@team_leader_required
def debug_template(request):
    """
    Диагностическая функция для анализа содержимого PPTX шаблона
    """
    try:
        template_path = os.path.join(settings.BASE_DIR, 'static', 'templates', 'certificate_template.pptx')
        prs = Presentation(template_path)
        
        result = []
        
        for i, slide in enumerate(prs.slides):
            slide_data = {
                "slide_number": i + 1,
                "shapes": []
            }
            
            for j, shape in enumerate(slide.shapes):
                shape_data = {
                    "shape_number": j + 1,
                    "shape_type": str(shape.shape_type),
                    "has_text": hasattr(shape, 'text_frame') and shape.has_text_frame,
                    "text_content": []
                }
                
                if hasattr(shape, 'text_frame') and shape.has_text_frame:
                    for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                        for r_idx, run in enumerate(paragraph.runs):
                            shape_data["text_content"].append({
                                "paragraph": p_idx + 1,
                                "run": r_idx + 1,
                                "text": run.text
                            })
                
                slide_data["shapes"].append(shape_data)
            
            result.append(slide_data)
        
        return JsonResponse({"template_analysis": result})
    
    except Exception as e:
        import traceback
        trace = traceback.format_exc()
        return JsonResponse({"error": str(e), "trace": trace}, status=400)

@team_leader_required
def all_scanners_list(request):
    """
    Отображает список всех сканеров с возможностью фильтрации и поиска
    """
    query = request.GET.get('q', '')
    filter_min_hours = request.GET.get('min_hours', '')
    filter_max_hours = request.GET.get('max_hours', '')
    page_number = request.GET.get('page', '1')
    
    # Создаем кеш-ключ на основе параметров запроса
    cache_key = f'scanners_list_{query}_{filter_min_hours}_{filter_max_hours}_{page_number}'
    cached_data = cache.get(cache_key)
    
    if cached_data:
        return render(request, 'core/all_scanners.html', cached_data)
    
    # Базовый запрос со всеми сканерами
    scanners = Scanner.objects.all()
    
    # Применяем фильтр по имени/email
    if query:
        scanners = scanners.filter(
            Q(first_name__icontains=query) | 
            Q(last_name__icontains=query) | 
            Q(email__icontains=query)
        )
    
    # Используем annotate для расчета суммарных часов вместо Python-циклов
    scanners = scanners.annotate(
        total_hours=Coalesce(Sum('eventparticipant__hours_awarded'), 0.0)
    )
    
    # Применяем фильтры по часам на уровне базы данных (учитываем общее количество часов)
    if filter_min_hours:
        min_hours = float(filter_min_hours)
        scanners = scanners.filter(Q(total_hours__gte=min_hours) | Q(total_certificate_hours__gte=min_hours) | 
                                  Q(total_hours__add=F('total_certificate_hours'))>=min_hours)
    
    if filter_max_hours:
        max_hours = float(filter_max_hours)
        scanners = scanners.filter(Q(total_hours__lte=max_hours) & Q(total_certificate_hours__lte=max_hours) & 
                                  Q(total_hours__add=F('total_certificate_hours'))<=max_hours)
    
    # Сортируем по убыванию общего количества часов
    scanners = scanners.annotate(
        all_hours=F('total_hours') + F('total_certificate_hours')
    ).order_by('-all_hours', 'last_name', 'first_name')
    
    # Пагинация
    paginator = Paginator(scanners, 20)  # 20 сканеров на страницу
    page_obj = paginator.get_page(page_number)
    
    # Предзагружаем данные о мероприятиях для сканеров на текущей странице
    scanner_ids = [scanner.id for scanner in page_obj.object_list]
    
    # Получаем все участия для сканеров на текущей странице одним запросом
    participations = EventParticipant.objects.filter(
        volunteer_id__in=scanner_ids
    ).select_related('event').order_by('-event__date')
    
    # Группируем участия по сканерам
    participation_by_scanner = {}
    for p in participations:
        if p.volunteer_id not in participation_by_scanner:
            participation_by_scanner[p.volunteer_id] = []
        
        participation_by_scanner[p.volunteer_id].append({
            'id': p.event.id,
            'name': p.event.name,
            'date': p.event.date.strftime('%d.%m.%Y'),
            'hours': p.hours_awarded
        })
    
    # Формируем итоговый список сканеров с их мероприятиями
    scanners_with_hours = []
    for scanner in page_obj.object_list:
        scanners_with_hours.append({
            'id': scanner.id,
            'first_name': scanner.first_name,
            'last_name': scanner.last_name,
            'email': scanner.email,
            'total_hours': scanner.total_hours,
            'total_certificate_hours': scanner.total_certificate_hours,
            'events': participation_by_scanner.get(scanner.id, [])
        })
    
    context = {
        'scanners': scanners_with_hours,
        'query': query,
        'min_hours': filter_min_hours,
        'max_hours': filter_max_hours,
        'page_obj': page_obj
    }
    
    # Кешируем результат на 5 минут
    cache.set(cache_key, context, 300)
    
    return render(request, 'core/all_scanners.html', context)

def create_certificate_pdf(name, hours, event_name=None, event_date=None, leader_name=None, period=None, events_list=None):
    from PIL import Image as PILImage, ImageDraw, ImageFont
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import landscape, A4
    import tempfile, os

    base_dir = settings.BASE_DIR
    bg_path = os.path.join(base_dir, 'static', 'templates', 'background.png')
    logo_path = os.path.join(base_dir, 'static', 'templates', 'image.png')
    stamp_path = os.path.join(base_dir, 'static', 'templates', 'stamp.png')
    impact_path = os.path.join(base_dir, 'static', 'fonts', 'IMPACT.TTF')
    sans_italic = os.path.join(base_dir, 'static', 'fonts', 'OpenSans-Italic.ttf')
    sans_bold_italic = os.path.join(base_dir, 'static', 'fonts', 'OpenSans-BoldItalic.ttf')

    # Увеличиваем разрешение для максимального качества
    scale = 2
    width, height = 1123 * scale, 794 * scale
    
    # Создаем полностью черный фон как базовый слой
    img = PILImage.new('RGB', (width, height), (0, 0, 0))
    
    # Загружаем и накладываем основной фон
    bg = PILImage.open(bg_path).convert('RGBA').resize((width, height))
    img.paste(bg, (0, 0), bg)
    
    # Создаем полупрозрачный черный слой вместо использования back_black.png
    # Создаем новое RGBA изображение с черным цветом и прозрачностью 60%
    overlay = PILImage.new('RGBA', (width, height), (0, 0, 0, 153))  # 153 = 60% непрозрачности
    
    # Накладываем полупрозрачный черный слой поверх основного фона
    img.paste((0, 0, 0), (0, 0, width, height), overlay)
    
    draw = ImageDraw.Draw(img)

    # Логотип в левый нижний угол (logo_scale = 0.9)
    logo = PILImage.open(logo_path).convert('RGBA')
    logo_w, logo_h = logo.size
    logo_scale = 0.9
    logo = logo.resize((int(logo_w * logo_scale), int(logo_h * logo_scale)), PILImage.LANCZOS)
    img.paste(logo, (40 * scale, height - logo.height - 40 * scale), logo)

    # Шрифты
    impact = ImageFont.truetype(impact_path, size=98)
    impact_name = ImageFont.truetype(impact_path, size=128)
    sans_bold_italic_f = ImageFont.truetype(sans_bold_italic, size=24)
    sans_italic_f = ImageFont.truetype(sans_italic, size=18)
    sans_italic_f_small = ImageFont.truetype(sans_italic, size=16)
    sans_bold_italic_f_small = ImageFont.truetype(sans_bold_italic, size=18)

    def get_text_size(text, font):
        bbox = font.getbbox(text)
        return bbox[2] - bbox[0], bbox[3] - bbox[1]

    # CERTIFICAT
    cert_text = "CERTIFICAT"
    impact_cert = ImageFont.truetype(impact_path, size=int(98.3 * scale))
    cert_w, cert_h = get_text_size(cert_text, impact_cert)
    cert_x = (width - cert_w - 80 * scale)
    cert_y = 70 * scale
    draw.text((cert_x, cert_y), cert_text, font=impact_cert, fill=(255,255,255,255))

    # OF APPRECIATION
    app_text = "OF APPRECIATION"
    app_font = ImageFont.truetype(sans_bold_italic, size=24 * scale)
    app_w, app_h = get_text_size(app_text, app_font)
    app_x = (width - app_w - 85 * scale)
    app_y = cert_y + cert_h + 40 * scale
    draw.text((app_x, app_y), app_text, font=app_font, fill=(255,255,255,255))

    # This certificate is presented to:
    pres_text = "This certificate is presented to:"
    pres_font = ImageFont.truetype(sans_italic, size=18 * scale)
    pres_w, pres_h = get_text_size(pres_text, pres_font)
    pres_x = (width - pres_w - 120 * scale)
    pres_y = app_y + app_h + 40 * scale
    draw.text((pres_x, pres_y), pres_text, font=pres_font, fill=(255,255,255,255))

    # NAME
    name_text = name
    name_w, name_h = get_text_size(name_text, impact_name)
    name_x = (width - name_w - 120 * scale)
    name_y = pres_y + pres_h + 10 * scale
    draw.text((name_x, name_y), name_text, font=impact_name, fill=(76,175,80,255))

    # Определяем годы для текста благодарности
    years_text = "2024 and 2025"  # По умолчанию
    
    # Если есть список мероприятий, извлекаем годы из них
    if events_list:
        years = set()
        for event in events_list:
            if 'date' in event:
                try:
                    # Извлекаем год из даты в формате DD.MM.YYYY
                    date_parts = event['date'].split('.')
                    if len(date_parts) == 3:
                        years.add(date_parts[2])  # Год - последняя часть
                except:
                    pass
        
        # Если нашли годы, формируем строку
        if years:
            sorted_years = sorted(years)
            if len(sorted_years) == 1:
                years_text = sorted_years[0]
            elif len(sorted_years) == 2:
                years_text = f"{sorted_years[0]} and {sorted_years[1]}"
            else:
                # Для трех и более лет используем перечисление с запятыми и "and" перед последним
                years_text = ", ".join(sorted_years[:-1]) + " and " + sorted_years[-1]
    
    # Если передан период, извлекаем годы из него
    elif period:
        try:
            # Предполагаем, что период имеет формат "DD.MM.YYYY - DD.MM.YYYY"
            period_parts = period.split(' - ')
            if len(period_parts) == 2:
                start_year = period_parts[0].split('.')[-1]
                end_year = period_parts[1].split('.')[-1]
                
                if start_year == end_year:
                    years_text = start_year
                else:
                    years_text = f"{start_year} and {end_year}"
        except:
            pass
    
    # Если передана одиночная дата мероприятия
    elif event_date:
        try:
            # Предполагаем формат даты "DD.MM.YYYY" или объект datetime
            if isinstance(event_date, str):
                year = event_date.split('.')[-1]
                years_text = year
            else:
                years_text = str(event_date.year)
        except:
            pass

    # Блок благодарности с динамическими годами
    lines = [
        "We, the Ticketon company, would like to sincerely express our gratitude and",
        "appreciation towards your incredible work and support in organizing",
        f"our events in {years_text}. You played important role in organization of",
        "each event. We hope to see you again in upcoming events!"
    ]
    main_font = ImageFont.truetype(sans_italic, size=16 * scale)
    main_width = 600 * scale
    main_x = width - main_width - 40 * scale
    main_y = name_y + name_h + 40 * scale
    line_h = main_font.getbbox('Ag')[3] - main_font.getbbox('Ag')[1] + 8 * scale
    for i, line in enumerate(lines):
        words = line.split()
        if len(words) == 1:
            draw.text((main_x, main_y + i * line_h), line, font=main_font, fill=(255,255,255,255))
            continue
        total_w = sum(get_text_size(w, main_font)[0] for w in words)
        space_w = int((main_width - total_w) / max(1, (len(words) - 1) * 2))
        x = main_x
        for j, word in enumerate(words):
            draw.text((x, main_y + i * line_h), word, font=main_font, fill=(255,255,255,255))
            w, _ = get_text_size(word, main_font)
            x += w + space_w

    # Рисуем стрелку с остриём только слева, справа — ровно
    arrow_w, arrow_height = 540 * scale, 110 * scale
    arrow = PILImage.new('RGBA', (arrow_w, arrow_height), (0,0,0,0))
    adraw = ImageDraw.Draw(arrow)
    points = [
        (0, arrow_height//2), (40 * scale, 0), (arrow_w, 0), (arrow_w, arrow_height), (40 * scale, arrow_height)
    ]
    adraw.polygon(points, fill=(76,175,80,255))

    margin = 18 * scale
    gap = 14 * scale
    gap_hours = 6 * scale
    section_w = (arrow_w - 2 * margin - 2 * gap) // 3
    center_y = arrow_height // 2

    # Шрифты для стрелки
    impact_25 = ImageFont.truetype(impact_path, size=25 * scale)
    impact_19_9 = ImageFont.truetype(impact_path, size=int(19.9 * scale))
    impact_18 = ImageFont.truetype(impact_path, size=18 * scale)

    # Левая секция: часы
    hours_text = f"{int(round(hours)):02d}"
    hw, hh = get_text_size(hours_text, impact_25)
    hlabel = "hours"
    hlw, hlh = get_text_size(hlabel, impact_18)
    line_w = section_w * 0.8
    total_h = hh + gap + hlh
    base_y = center_y - total_h // 2 - 10 * scale
    hours_x = margin + (section_w - hw) // 2
    hours_y = base_y
    adraw.text((hours_x, hours_y), hours_text, font=impact_25, fill=(255,255,255,255))
    line_y = hours_y + hh + gap
    line_x1 = margin + (section_w - line_w) // 2
    line_x2 = line_x1 + line_w
    adraw.line([(line_x1, line_y), (line_x2, line_y)], fill=(255,255,255,255), width=4 * scale)
    hlabel_x = margin + (section_w - hlw) // 2
    hlabel_y = line_y + gap_hours
    adraw.text((hlabel_x, hlabel_y), hlabel, font=impact_18, fill=(255,255,255,255))

    # Центральная секция: печать (штамп чуть выше, овальный)
    if os.path.exists(stamp_path):
        stamp = PILImage.open(stamp_path).convert('RGBA')
        stamp_w = int(section_w)
        stamp_h = int(arrow_height * 0.95)
        stamp = stamp.resize((stamp_w, stamp_h), PILImage.LANCZOS)
        stamp_x = margin + section_w + gap
        stamp_y = center_y - stamp_h // 2
        arrow.paste(stamp, (stamp_x, stamp_y), stamp)

    # Правая секция: крупный текст, увеличиваем размер имени директора и слова "director"
    sign_text = "Torgunakova V. K."
    dir_text = "director"
    max_width = section_w - 2 * int(5 * scale)  # Уменьшаем внутренний отступ для большего текста
    
    # Увеличиваем базовый размер шрифта для имени директора и слова "director"
    base_sign_font_size = int(25 * 1.5 * scale)  # Было 20.9 * 1.5
    base_dir_font_size = int(22 * 1.5 * scale)   # Было 19 * 1.5
    
    min_font_size = int(15 * scale)  # Увеличиваем минимальный размер шрифта
    sign_font_size = base_sign_font_size
    dir_font_size = base_dir_font_size
    
    while True:
        sign_font = ImageFont.truetype(impact_path, size=sign_font_size)
        dir_font = ImageFont.truetype(impact_path, size=dir_font_size)
        sign_w, sign_h = get_text_size(sign_text, sign_font)
        dir_w, dir_h = get_text_size(dir_text, dir_font)
        if sign_w <= max_width and dir_w <= max_width:
            break
        sign_font_size -= 2
        dir_font_size -= 2
        if sign_font_size < min_font_size or dir_font_size < min_font_size:
            sign_font_size = dir_font_size = min_font_size
            sign_font = ImageFont.truetype(impact_path, size=sign_font_size)
            dir_font = ImageFont.truetype(impact_path, size=dir_font_size)
            sign_w, sign_h = get_text_size(sign_text, sign_font)
            dir_w, dir_h = get_text_size(dir_text, dir_font)
            break
    
    right_section_x = margin + 2 * section_w + 2 * gap
    right_section_y = margin
    sign_x = right_section_x + (section_w - sign_w) // 2
    sign_y = right_section_y + 10 * scale
    dir_x = right_section_x + (section_w - dir_w) // 2
    dir_y = sign_y + sign_h + int(0.2 * sign_h)  # Уменьшаем расстояние между именем и должностью
    
    # Рисуем имя и должность директора
    adraw.text((sign_x, sign_y), sign_text, font=sign_font, fill=(0,0,0,255))
    adraw.text((dir_x, dir_y), dir_text, font=dir_font, fill=(255,255,255,255))

    arrow_x = width - arrow_w
    arrow_y = height - arrow_height - 60 * scale
    img.paste(arrow, (arrow_x, arrow_y), arrow)

    temp_dir = tempfile.mkdtemp()
    temp_img_path = os.path.join(temp_dir, 'cert.png')
    img.save(temp_img_path, 'PNG')

    # Использование ReportLab для создания PDF без артефактов
    pdf_width, pdf_height = 1123, 794  # Размеры PDF в точках
    temp_pdf_path = os.path.join(temp_dir, 'certificate.pdf')
    
    # Создаем PDF с чистым черным фоном, без артефактов
    c = canvas.Canvas(temp_pdf_path, pagesize=(pdf_width, pdf_height))
    
    # Заливаем весь PDF черным цветом (без границ)
    c.setFillColor((0, 0, 0))
    c.rect(0, 0, pdf_width, pdf_height, fill=1, stroke=0)
    
    # Добавляем изображение сертификата (без верхней черной линии)
    c.drawImage(temp_img_path, 0, 0, width=pdf_width, height=pdf_height)
    
    c.save()

    with open(temp_pdf_path, 'rb') as f:
        pdf_data = f.read()
    
    # Очистка временных файлов
    os.remove(temp_img_path)
    os.remove(temp_pdf_path)
    os.rmdir(temp_dir)
    
    return pdf_data

def convert_pptx_to_pdf(pptx_path):
    """
    Конвертирует PPTX в PDF используя PowerPoint через COM-интерфейс (Windows)
    """
    pdf_path = pptx_path.replace('.pptx', '.pdf')
    
    # Определяем константы для PowerPoint
    try:
        import comtypes.client
        import pythoncom
        
        # Инициализация COM
        pythoncom.CoInitialize()
        
        # Форматы для экспорта (константы для PowerPoint)
        ppSaveAsPDF = 32  # Формат PDF
        
        # Создаем экземпляр PowerPoint
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = True  # PowerPoint должен быть видимым из-за ограничений безопасности
        
        # Полный абсолютный путь к файлам
        abs_pptx_path = os.path.abspath(pptx_path)
        abs_pdf_path = os.path.abspath(pdf_path)
        
        print(f"Открываем презентацию: {abs_pptx_path}")
        
        # Открываем презентацию
        presentation = powerpoint.Presentations.Open(abs_pptx_path)
        
        # Небольшая пауза для загрузки презентации
        time.sleep(2)
        
        try:
            print(f"Сохраняем как PDF: {abs_pdf_path}")
            # Сохраняем как PDF
            presentation.SaveAs(abs_pdf_path, ppSaveAsPDF)
            
            # Пауза для сохранения PDF
            time.sleep(2)
            
            # Закрываем презентацию
            presentation.Close()
        except Exception as e:
            print(f"Ошибка при сохранении PDF: {e}")
        finally:
            # Закрываем PowerPoint
            powerpoint.Quit()
        
        # Освобождаем COM-объекты
        pythoncom.CoUninitialize()
        
        # Проверяем, создан ли PDF-файл
        if os.path.exists(pdf_path):
            print(f"PDF файл успешно создан: {pdf_path}")
            with open(pdf_path, 'rb') as f:
                pdf_data = f.read()
            return pdf_data
        else:
            print(f"PDF-файл не создан по пути {pdf_path}")
            return None
    except Exception as e:
        import traceback
        print(f"Ошибка при конвертации PPTX в PDF: {e}")
        print(traceback.format_exc())
        return None

@team_leader_required
def event_delete(request, event_id):
    event = get_object_or_404(Event, id=event_id)
    if not request.user.is_staff:
        messages.error(request, 'Удаление доступно только администраторам.')
        return redirect('event_edit', event_id=event.id)
    if request.method == 'POST':
        event.delete()
        messages.success(request, 'Мероприятие успешно удалено.')
        return redirect('events')
    return render(request, 'core/event_confirm_delete.html', {'event': event})

# Функция для очистки кеша связанного со сканерами
def clear_scanners_cache():
    """Очищает все кеши, связанные со сканерами"""
    # Очищаем кеш по префиксу
    cache.clear()

@team_leader_required
def purge_settings(request):
    """Страница настроек автоматического удаления мероприятий"""
    purge_config, created = PurgeSettings.objects.get_or_create(
        defaults={
            'purge_date': timezone.datetime(timezone.now().year, 9, 1).date(),
            'notification_days_before': 7,
            'active': True,
            'updated_by': request.user
        }
    )
    
    if request.method == 'POST':
        action = request.POST.get('action')
        
        # Обновление настроек
        if action == 'update':
            try:
                # Получаем и валидируем дату удаления
                purge_month = int(request.POST.get('purge_month', 9))
                purge_day = int(request.POST.get('purge_day', 1))
                
                # Проверяем валидность даты
                if purge_month < 1 or purge_month > 12 or purge_day < 1 or purge_day > 31:
                    messages.error(request, 'Некорректная дата удаления')
                    return redirect('purge_settings')
                
                # Создаем дату удаления
                purge_date = timezone.datetime(timezone.now().year, purge_month, purge_day).date()
                
                # Получаем и валидируем дни для уведомлений
                notification_days = int(request.POST.get('notification_days', 7))
                if notification_days < 1 or notification_days > 30:
                    messages.error(request, 'Количество дней для уведомления должно быть от 1 до 30')
                    return redirect('purge_settings')
                
                # Получаем статус активности
                active = request.POST.get('active') == 'on'
                
                # Обновляем настройки
                purge_config.purge_date = purge_date
                purge_config.notification_days_before = notification_days
                purge_config.active = active
                purge_config.updated_by = request.user
                purge_config.updated_at = timezone.now()
                purge_config.save()
                
                messages.success(request, 'Настройки успешно обновлены')
            except Exception as e:
                messages.error(request, f'Ошибка при обновлении настроек: {str(e)}')
        
        return redirect('purge_settings')
    
    # Рассчитываем дату следующего удаления
    current_date = timezone.now().date()
    purge_month_day = purge_config.purge_date.strftime('%m-%d')
    next_purge_year = current_date.year
    
    # Если текущая дата после даты удаления в этом году, то следующее удаление в следующем году
    next_purge_date_str = f"{next_purge_year}-{purge_month_day}"
    next_purge_date = datetime.strptime(next_purge_date_str, '%Y-%m-%d').date()
    
    if current_date > next_purge_date:
        next_purge_year += 1
        next_purge_date_str = f"{next_purge_year}-{purge_month_day}"
        next_purge_date = datetime.strptime(next_purge_date_str, '%Y-%m-%d').date()
    
    # Рассчитываем дату уведомления
    notification_date = next_purge_date - timedelta(days=purge_config.notification_days_before)
    
    # Получаем количество мероприятий, которые будут удалены
    one_year_ago = current_date - timedelta(days=365)
    events_to_delete_count = Event.objects.filter(date__lt=one_year_ago).count()
    
    # Получаем последние логи уведомлений
    recent_logs = NotificationLog.objects.filter(is_test=False).order_by('-sent_at')[:5]
    
    return render(request, 'core/purge_settings.html', {
        'settings': purge_config,
        'next_purge_date': next_purge_date,
        'notification_date': notification_date,
        'events_to_delete_count': events_to_delete_count,
        'recent_logs': recent_logs
    })

@admin_required
def notification_logs(request):
    """Страница с логами отправленных уведомлений"""
    logs = NotificationLog.objects.all()
    
    # Фильтрация по типу (тестовое/системное)
    is_test = request.GET.get('is_test')
    if is_test == 'true':
        logs = logs.filter(is_test=True)
    elif is_test == 'false':
        logs = logs.filter(is_test=False)
    
    # Фильтрация по способу доставки (email/telegram)
    notification_type = request.GET.get('notification_type')
    if notification_type:
        logs = logs.filter(notification_type=notification_type)
    
    # Поиск по получателю (email или telegram_id)
    recipient = request.GET.get('recipient', '').strip()
    if recipient:
        logs = logs.filter(
            Q(recipient_email__icontains=recipient) | 
            Q(recipient_telegram_id__icontains=recipient)
        )
    
    # Пагинация
    paginator = Paginator(logs, 20)
    page_number = request.GET.get('page')
    page_obj = paginator.get_page(page_number)
    
    context = {
        'logs': page_obj.object_list,
        'page_obj': page_obj,
        'is_test': is_test,
        'recipient': recipient,
        'notification_type': notification_type
    }
    
    return render(request, 'core/notification_logs.html', context)

# Обновление команды для включения пользовательских настроек
def update_purge_command():
    """Обновляет команду purge_events с учетом пользовательских настроек"""
    try:
        purge_config = PurgeSettings.objects.first()
        if not purge_config:
            return
        
        # Обновляем даты в скрипте schedule_tasks.py
        with open('schedule_tasks.py', 'r') as f:
            content = f.read()
        
        # Получаем месяц и день из настроек
        month = purge_config.purge_date.month
        day = purge_config.purge_date.day
        
        # Заменяем дату в скрипте
        next_year = timezone.now().year + 1
        updated_content = re.sub(
            r'start_date = datetime\(\d+, \d+, \d+, \d+, \d+, \d+\)',
            f'start_date = datetime({next_year}, {month}, {day}, 2, 0, 0)',
            content
        )
        
        with open('schedule_tasks.py', 'w') as f:
            f.write(updated_content)
        
        return True
    except Exception as e:
        print(f"Ошибка при обновлении команды: {e}")
        return False

@team_leader_required
def send_test_notification(request):
    """Отдельная страница для отправки тестовых уведомлений"""
    if request.method == 'POST':
        # Получаем настройки удаления
        purge_settings, created = PurgeSettings.objects.get_or_create(
            defaults={
                'purge_date': timezone.datetime(timezone.now().year, 7, 13).date(),  # 13 июля
                'notification_days_before': 7,
                'active': True,
                'updated_by': request.user
            }
        )
        
        # Рассчитываем дату следующего удаления
        current_date = timezone.now().date()
        purge_month_day = purge_settings.purge_date.strftime('%m-%d')
        next_purge_year = current_date.year
        
        # Если текущая дата после даты удаления в этом году, то следующее удаление в следующем году
        next_purge_date_str = f"{next_purge_year}-{purge_month_day}"
        next_purge_date = datetime.strptime(next_purge_date_str, '%Y-%m-%d').date()
        
        if current_date > next_purge_date:
            next_purge_year += 1
            next_purge_date_str = f"{next_purge_year}-{purge_month_day}"
            next_purge_date = datetime.strptime(next_purge_date_str, '%Y-%m-%d').date()
        
        # Проверяем, какой тип уведомления выбран
        notification_type = request.POST.get('notification_type', 'telegram')
        
        # Проверяем, есть ли указанный получатель для тестирования
        test_recipient = request.POST.get('test_recipient', '').strip()
        
        # Получаем всех тимлидеров или конкретного получателя
        if test_recipient:
            # Отправляем одному указанному адресату
            if notification_type == 'email':
                recipients = [{'email': test_recipient, 'name': 'Тимлидер'}]
            else:
                recipients = [{'telegram_id': test_recipient, 'name': 'Тимлидер'}]
        else:
            # Отправляем всем тимлидерам
            team_leaders = User.objects.filter(groups__name='Тимлидеры')
            
            if not team_leaders.exists():
                messages.warning(request, 'Нет тимлидеров для отправки уведомлений')
                return redirect('send_test_notification')
            
            recipients = []
            for leader in team_leaders:
                recipient_data = {
                    'name': f"{leader.first_name} {leader.last_name}"
                }
                
                # Ищем Telegram ID для тимлидера
                if notification_type == 'telegram':
                    try:
                        team_leader = TeamLeader.objects.filter(
                            first_name=leader.first_name,
                            last_name=leader.last_name
                        ).first()
                        if team_leader and team_leader.telegram_id:
                            recipient_data['telegram_id'] = team_leader.telegram_id
                            recipients.append(recipient_data)
                    except Exception as e:
                        messages.error(request, f'Ошибка при поиске Telegram ID для {recipient_data["name"]}: {str(e)}')
                else:
                    if leader.email:
                        recipient_data['email'] = leader.email
                        recipients.append(recipient_data)
        
        # Проверяем, есть ли получатели
        if not recipients:
            if notification_type == 'telegram':
                messages.error(request, 'Не найдены Telegram ID для отправки уведомлений')
            else:
                messages.error(request, 'Не найдены email-адреса для отправки уведомлений')
            return redirect('send_test_notification')
        
        sent_count = 0
        from core.utils import send_telegram_message
        
        for recipient in recipients:
            # Формируем текст тестового уведомления
            message = f"""
            Здравствуйте, {recipient['name']}!

            Это тестовое уведомление системы автоматического удаления мероприятий.
            
            Следующее удаление мероприятий запланировано на {next_purge_date.strftime('%d.%m.%Y')}.
            Будут удалены мероприятия старше одного года.
            
            Это сообщение отправлено администратором {request.user.first_name} {request.user.last_name} для проверки работы системы уведомлений.

            С уважением,
            Команда Ticketon
            """
            
            # Отправляем уведомление
            try:
                if notification_type == 'telegram' and 'telegram_id' in recipient:
                    # Отправляем через Telegram
                    success = send_telegram_message(message, [recipient['telegram_id']])
                    
                    if success:
                        # Логируем отправку
                        NotificationLog.objects.create(
                            sent_by=request.user,
                            recipient_telegram_id=recipient['telegram_id'],
                            message=message,
                            is_test=True,
                            notification_type='telegram'
                        )
                        
                        sent_count += 1
                    else:
                        messages.error(request, f'Не удалось отправить уведомление в Telegram для {recipient["name"]}')
                
                elif notification_type == 'email' and 'email' in recipient:
                    # Отправляем через email
                    send_mail(
                        subject='Тестовое уведомление: система удаления мероприятий',
                        message=message,
                        from_email=settings.DEFAULT_FROM_EMAIL,
                        recipient_list=[recipient['email']],
                        fail_silently=False,
                    )
                    
                    # Логируем отправку
                    NotificationLog.objects.create(
                        sent_by=request.user,
                        recipient_email=recipient['email'],
                        message=message,
                        is_test=True,
                        notification_type='email'
                    )
                    
                    sent_count += 1
            except Exception as e:
                recipient_id = recipient.get('telegram_id', recipient.get('email', 'неизвестный получатель'))
                messages.error(request, f'Ошибка при отправке уведомления на {recipient_id}: {str(e)}')
        
        if sent_count > 0:
            if test_recipient:
                messages.success(request, f'Тестовое уведомление отправлено на {test_recipient}')
            else:
                messages.success(request, f'Тестовые уведомления отправлены {sent_count} тимлидерам')
        
        return redirect('events')
    
    # Получаем последние отправленные уведомления
    recent_notifications = NotificationLog.objects.filter(is_test=True)[:5]
    
    return render(request, 'core/send_test_notification.html', {
        'recent_notifications': recent_notifications
    })
